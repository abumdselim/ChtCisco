"""
hydronet_generator.py
=====================
Generates HydroNet_Chittagong_Resilience_Final.pptx — a 14-slide, production-grade
university project presentation on:
  "CHITTAGONG HYDROLOGICAL RESILIENCE — A Multi-Zone IoT Network for
   Flash Flood & Tidal Surge Mitigation"

Dependencies: python-pptx, networkx, matplotlib, Pillow, numpy
Run:  python hydronet_generator.py
"""

import io
import math
import os
import textwrap

import matplotlib
matplotlib.use("Agg")
import matplotlib.patches as mpatches
import matplotlib.pyplot as plt
import networkx as nx
import numpy as np
from PIL import Image, ImageDraw, ImageFont

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

# ─────────────────────────────────────────────────────────────────────────────
# COLOUR PALETTE
# ─────────────────────────────────────────────────────────────────────────────
PRIMARY_BLUE  = RGBColor(0x1a, 0x23, 0x32)
TEAL_ACCENT   = RGBColor(0x00, 0xd4, 0xff)
WHITE_CLEAN   = RGBColor(0xff, 0xff, 0xff)
GREEN_ACCENT  = RGBColor(0x2e, 0xcc, 0x71)
DARK_SLATE    = RGBColor(0x0f, 0x14, 0x19)
LIGHT_GRAY    = RGBColor(0xe8, 0xea, 0xed)
ORANGE_ALERT  = RGBColor(0xff, 0x6b, 0x35)
PURPLE_ZONE   = RGBColor(0x9b, 0x59, 0xb6)

# hex helpers for matplotlib
_HEX = {
    "navy":   "#1a2332",
    "teal":   "#00d4ff",
    "white":  "#ffffff",
    "green":  "#2ecc71",
    "slate":  "#0f1419",
    "lgray":  "#e8eaed",
    "orange": "#ff6b35",
    "purple": "#9b59b6",
}

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ─────────────────────────────────────────────────────────────────────────────
# HELPERS
# ─────────────────────────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill: RGBColor | None = None,
             line: RGBColor | None = None, line_w: float = 1.0):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_name="Calibri",
                font_size=12, bold=False, italic=False, color=None,
                align=PP_ALIGN.LEFT, word_wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return txb


def add_para(tf, text, font_name="Calibri", font_size=11, bold=False,
             italic=False, color=None, align=PP_ALIGN.LEFT, space_before=0):
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return p


def fig_to_stream(fig, dpi=150):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=dpi, bbox_inches="tight")
    buf.seek(0)
    plt.close(fig)
    return buf


def add_underline_bar(slide, left, top, width, color: RGBColor, height=Pt(2)):
    return add_rect(slide, left, top, width, height, fill=color)


def set_slide_background(slide, color: RGBColor):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


# ─────────────────────────────────────────────────────────────────────────────
# DIAGRAM GENERATORS
# ─────────────────────────────────────────────────────────────────────────────

def make_arch_overview_diagram() -> io.BytesIO:
    """Slide 5 – three-tier architecture overview."""
    fig, ax = plt.subplots(figsize=(14, 8))
    fig.patch.set_facecolor(_HEX["lgray"])
    ax.set_facecolor(_HEX["lgray"])
    ax.axis("off")

    def box(ax, x, y, w, h, txt, bg, fg="white", border=_HEX["teal"], lw=2, fs=9):
        rect = mpatches.FancyBboxPatch(
            (x - w / 2, y - h / 2), w, h,
            boxstyle="round,pad=0.02",
            facecolor=bg, edgecolor=border, linewidth=lw
        )
        ax.add_patch(rect)
        ax.text(x, y, txt, ha="center", va="center", color=fg,
                fontsize=fs, fontweight="bold", wrap=True,
                multialignment="center")

    def arrow(ax, x1, y1, x2, y2, color=_HEX["teal"], lw=2, style="->"):
        ax.annotate("", xy=(x2, y2), xytext=(x1, y1),
                    arrowprops=dict(arrowstyle=style, color=color, lw=lw))

    # ── Tier labels ──────────────────────────────────────────────────────────
    for label, ypos in [("TIER 1 — Core Layer", 9),
                         ("TIER 2 — Distribution Layer", 6.2),
                         ("TIER 3 — Access & IoT Layer", 3.2)]:
        ax.text(0.3, ypos, label, ha="left", va="center",
                color=_HEX["navy"], fontsize=9, fontstyle="italic",
                fontweight="bold")

    # ── OSPF area 0 band ─────────────────────────────────────────────────────
    ax.add_patch(mpatches.FancyBboxPatch((0.5, 8), 13, 1.3,
                 boxstyle="round,pad=0.05",
                 facecolor="#d6eaf8", edgecolor=_HEX["teal"], linewidth=1.5,
                 alpha=0.5))
    ax.text(7, 8.65, "OSPF Area 0  (Backbone — 10.0.0.0/16)",
            ha="center", color=_HEX["navy"], fontsize=9, fontstyle="italic")

    # Core routers
    box(ax, 4.5, 8.65, 2.6, 0.85,
        "Core Router 1\n(3945 ISR) | 10.0.0.1", _HEX["navy"], lw=3, fs=9)
    box(ax, 9.5, 8.65, 2.6, 0.85,
        "Core Router 2\n(3945 ISR) | 10.0.1.1", _HEX["navy"], lw=3, fs=9)
    # Redundancy double-line
    ax.annotate("", xy=(7.2, 8.65), xytext=(5.8, 8.65),
                arrowprops=dict(arrowstyle="<->", color=_HEX["teal"], lw=2.5))
    ax.text(7, 9.1, "Redundancy", ha="center", fontsize=7.5, color=_HEX["teal"])

    # Core switch (server pool)
    box(ax, 7, 7.3, 2.4, 0.7,
        "Core Switch\nCatalyst 3650 | 10.0.0.10", _HEX["teal"],
        fg=_HEX["slate"], border=_HEX["navy"], fs=8)
    arrow(ax, 7, 7.65, 7, 8.2)

    # ── Zone distribution routers ─────────────────────────────────────────────
    zone_xs = [2, 5, 9, 12]
    zone_labels = [
        "Z1 Router\nCisco 2901\n10.1.0.1\nOSPF Area 1",
        "Z2 Router\nCisco 2901\n10.2.0.1\nOSPF Area 2",
        "Z3 Router\nCisco 2901\n10.3.0.1\nOSPF Area 3",
        "Z4 Router\nCisco 2901\n10.4.0.1\nOSPF Area 4",
    ]
    zone_bg = [_HEX["orange"], _HEX["teal"], _HEX["purple"], _HEX["green"]]
    for x, lbl, bg in zip(zone_xs, zone_labels, zone_bg):
        box(ax, x, 6.2, 2.2, 1.0, lbl, bg, fs=8)
        arrow(ax, x, 6.7, x if x not in (2, 12) else (4.5 if x == 2 else 9.5),
              8.2, color=bg, lw=1.8)

    # ── Access layer ─────────────────────────────────────────────────────────
    iot_colors = [_HEX["orange"], _HEX["teal"], _HEX["purple"], _HEX["green"]]
    for xi, (x, bg) in enumerate(zip(zone_xs, iot_colors)):
        # Switch
        box(ax, x, 4.65, 1.8, 0.55,
            f"2960-X Switch\n10.{xi+1}.0.5", "#2c3e50",
            fg="white", border=bg, fs=7.5)
        arrow(ax, x, 4.93, x, 5.7, color=bg, lw=1.5)
        # AP
        box(ax, x, 3.9, 1.8, 0.55,
            f"AP Zone-{xi+1}\nSSID: Z{xi+1}-HydroNet", "#8e44ad",
            fg="white", border=bg, fs=7.5)
        arrow(ax, x, 4.18, x, 4.38, color=bg, lw=1.5, style="->")
        # IoT circles
        for j in range(5):
            cx = x - 0.9 + j * 0.45
            cy = 3.05
            circ = plt.Circle((cx, cy), 0.18, color=_HEX["green"],
                               ec=_HEX["slate"], linewidth=0.8)
            ax.add_patch(circ)
        ax.text(x, 2.6, "5× IoT Sensors", ha="center", fontsize=7,
                color=_HEX["slate"])
        arrow(ax, x, 3.45, x, 3.62, color=_HEX["green"], lw=1.2)

    # ── Servers (bottom-centre) ───────────────────────────────────────────────
    box(ax, 7, 1.4, 3.8, 0.9,
        "DHCP | DNS | Web | Email | Syslog\n10.0.100.1 – 10.0.100.5  (VLAN 100)",
        _HEX["lgray"], fg=_HEX["slate"], border=_HEX["navy"], fs=8)
    arrow(ax, 7, 1.85, 7, 6.9, color=_HEX["navy"], lw=1.5)

    # ── Legend ────────────────────────────────────────────────────────────────
    legend_items = [
        (mpatches.Patch(color=_HEX["navy"]), "Core Router"),
        (mpatches.Patch(color=_HEX["teal"]), "Zone Router / Switch"),
        (mpatches.Patch(color=_HEX["green"]), "IoT Sensor"),
        (mpatches.Patch(color=_HEX["lgray"], ec=_HEX["navy"]), "Server"),
    ]
    ax.legend(
        [h for h, _ in legend_items],
        [l for _, l in legend_items],
        loc="lower right", fontsize=8,
        facecolor=_HEX["white"], edgecolor=_HEX["navy"]
    )

    ax.set_xlim(0, 14)
    ax.set_ylim(0.8, 10)
    fig.tight_layout()
    return fig_to_stream(fig, dpi=150)


def make_full_topology_diagram() -> io.BytesIO:
    """Slide 6 – full hierarchical network topology using networkx."""
    G = nx.DiGraph()

    nodes = {
        # Core
        "Core-R1": {"label": "Core-R1\n10.0.0.1", "color": _HEX["navy"],  "size": 2400, "layer": 4},
        "Core-R2": {"label": "Core-R2\n10.0.1.1", "color": _HEX["navy"],  "size": 2400, "layer": 4},
        "Core-SW": {"label": "Core-SW\n10.0.0.10","color": _HEX["teal"],  "size": 1800, "layer": 3},
        # Servers
        "Servers": {"label": "DHCP|DNS\nWeb|Mail|Syslog", "color": "#bdc3c7", "size": 2000, "layer": 2},
        # ISP
        "ISP":     {"label": "ISP\n172.16.0.1",    "color": _HEX["orange"], "size": 1600, "layer": 4},
        # Zone 1
        "Z1-R":  {"label": "Z1-Router\n10.1.0.1", "color": _HEX["orange"], "size": 1600, "layer": 3},
        "Z1-SW": {"label": "Z1-Switch\n10.1.0.5", "color": "#e67e22",       "size": 1200, "layer": 2},
        "Z1-AP": {"label": "Z1-AP",               "color": "#f39c12",       "size": 900,  "layer": 1},
        "Z1-WL": {"label": "WL-Z1\n10.1.1.10",   "color": _HEX["green"],  "size": 700,  "layer": 0},
        "Z1-RF": {"label": "Rain-Z1\n10.1.1.20",  "color": _HEX["green"],  "size": 700,  "layer": 0},
        "Z1-TS": {"label": "Tidal-Z1\n10.1.1.30", "color": _HEX["green"],  "size": 700,  "layer": 0},
        "Z1-TM": {"label": "Temp-Z1\n10.1.1.40",  "color": _HEX["green"],  "size": 700,  "layer": 0},
        "Z1-CA": {"label": "Cam-Z1\n10.1.1.50",   "color": _HEX["green"],  "size": 700,  "layer": 0},
        # Zone 2
        "Z2-R":  {"label": "Z2-Router\n10.2.0.1", "color": _HEX["teal"],  "size": 1600, "layer": 3},
        "Z2-SW": {"label": "Z2-Switch\n10.2.0.5", "color": "#1abc9c",     "size": 1200, "layer": 2},
        "Z2-AP": {"label": "Z2-AP",               "color": "#16a085",     "size": 900,  "layer": 1},
        "Z2-WL": {"label": "WL-Z2\n10.2.1.10",   "color": _HEX["green"], "size": 700,  "layer": 0},
        "Z2-RF": {"label": "Rain-Z2\n10.2.1.20",  "color": _HEX["green"], "size": 700,  "layer": 0},
        "Z2-HM": {"label": "Hum-Z2\n10.2.1.30",  "color": _HEX["green"], "size": 700,  "layer": 0},
        "Z2-TM": {"label": "Temp-Z2\n10.2.1.40",  "color": _HEX["green"], "size": 700,  "layer": 0},
        "Z2-CA": {"label": "Cam-Z2\n10.2.1.50",   "color": _HEX["green"], "size": 700,  "layer": 0},
        # Zone 3
        "Z3-R":  {"label": "Z3-Router\n10.3.0.1", "color": _HEX["purple"],"size": 1600, "layer": 3},
        "Z3-SW": {"label": "Z3-Switch\n10.3.0.5", "color": "#8e44ad",     "size": 1200, "layer": 2},
        "Z3-AP": {"label": "Z3-AP",               "color": "#9b59b6",     "size": 900,  "layer": 1},
        "Z3-WL": {"label": "WL-Z3\n10.3.1.10",   "color": _HEX["green"], "size": 700,  "layer": 0},
        "Z3-TS": {"label": "Tidal-Z3\n10.3.1.20", "color": _HEX["green"], "size": 700,  "layer": 0},
        "Z3-TM": {"label": "Temp-Z3\n10.3.1.30",  "color": _HEX["green"], "size": 700,  "layer": 0},
        "Z3-HM": {"label": "Hum-Z3\n10.3.1.40",  "color": _HEX["green"], "size": 700,  "layer": 0},
        "Z3-RF": {"label": "Rain-Z3\n10.3.1.50",  "color": _HEX["green"], "size": 700,  "layer": 0},
        # Zone 4
        "Z4-R":  {"label": "Z4-Router\n10.4.0.1", "color": _HEX["green"], "size": 1600, "layer": 3},
        "Z4-SW": {"label": "Z4-Switch\n10.4.0.5", "color": "#27ae60",     "size": 1200, "layer": 2},
        "Z4-AP": {"label": "Z4-AP",               "color": "#2ecc71",     "size": 900,  "layer": 1},
        "Z4-WL": {"label": "WL-Z4\n10.4.1.10",   "color": "#a9dfbf",     "size": 700,  "layer": 0},
        "Z4-RF": {"label": "Rain-Z4\n10.4.1.20",  "color": "#a9dfbf",     "size": 700,  "layer": 0},
        "Z4-HM": {"label": "Hum-Z4\n10.4.1.30",  "color": "#a9dfbf",     "size": 700,  "layer": 0},
        "Z4-TM": {"label": "Temp-Z4\n10.4.1.40",  "color": "#a9dfbf",     "size": 700,  "layer": 0},
        "Z4-CA": {"label": "Cam-Z4\n10.4.1.50",   "color": "#a9dfbf",     "size": 700,  "layer": 0},
    }

    for n, d in nodes.items():
        G.add_node(n, **d)

    edges = [
        # Core redundancy
        ("Core-R1", "Core-R2", {"style": "solid", "color": _HEX["teal"],  "width": 3}),
        ("Core-R1", "Core-SW", {"style": "solid", "color": _HEX["teal"],  "width": 2}),
        ("Core-R2", "Core-SW", {"style": "solid", "color": _HEX["teal"],  "width": 2}),
        # Core to servers
        ("Core-SW", "Servers", {"style": "solid", "color": _HEX["navy"],  "width": 2}),
        # Core to ISP
        ("Core-R1", "ISP",     {"style": "dashed","color": _HEX["orange"],"width": 1.5}),
        # Core to zone routers
        ("Core-R1", "Z1-R",   {"style": "solid", "color": _HEX["orange"],"width": 2}),
        ("Core-R1", "Z2-R",   {"style": "solid", "color": _HEX["teal"],  "width": 2}),
        ("Core-R2", "Z3-R",   {"style": "solid", "color": _HEX["purple"],"width": 2}),
        ("Core-R2", "Z4-R",   {"style": "solid", "color": _HEX["green"], "width": 2}),
        # Zone 1
        ("Z1-R","Z1-SW",{"style":"solid", "color":_HEX["orange"],"width":1.8}),
        ("Z1-SW","Z1-AP",{"style":"dashed","color":_HEX["orange"],"width":1.2}),
        ("Z1-AP","Z1-WL",{"style":"dotted","color":_HEX["green"],"width":1}),
        ("Z1-AP","Z1-RF",{"style":"dotted","color":_HEX["green"],"width":1}),
        ("Z1-AP","Z1-TS",{"style":"dotted","color":_HEX["green"],"width":1}),
        ("Z1-AP","Z1-TM",{"style":"dotted","color":_HEX["green"],"width":1}),
        ("Z1-AP","Z1-CA",{"style":"dotted","color":_HEX["green"],"width":1}),
        # Zone 2
        ("Z2-R","Z2-SW",{"style":"solid", "color":_HEX["teal"],  "width":1.8}),
        ("Z2-SW","Z2-AP",{"style":"dashed","color":_HEX["teal"],  "width":1.2}),
        ("Z2-AP","Z2-WL",{"style":"dotted","color":_HEX["green"],"width":1}),
        ("Z2-AP","Z2-RF",{"style":"dotted","color":_HEX["green"],"width":1}),
        ("Z2-AP","Z2-HM",{"style":"dotted","color":_HEX["green"],"width":1}),
        ("Z2-AP","Z2-TM",{"style":"dotted","color":_HEX["green"],"width":1}),
        ("Z2-AP","Z2-CA",{"style":"dotted","color":_HEX["green"],"width":1}),
        # Zone 3
        ("Z3-R","Z3-SW",{"style":"solid", "color":_HEX["purple"],"width":1.8}),
        ("Z3-SW","Z3-AP",{"style":"dashed","color":_HEX["purple"],"width":1.2}),
        ("Z3-AP","Z3-WL",{"style":"dotted","color":_HEX["green"],"width":1}),
        ("Z3-AP","Z3-TS",{"style":"dotted","color":_HEX["green"],"width":1}),
        ("Z3-AP","Z3-TM",{"style":"dotted","color":_HEX["green"],"width":1}),
        ("Z3-AP","Z3-HM",{"style":"dotted","color":_HEX["green"],"width":1}),
        ("Z3-AP","Z3-RF",{"style":"dotted","color":_HEX["green"],"width":1}),
        # Zone 4
        ("Z4-R","Z4-SW",{"style":"solid", "color":_HEX["green"], "width":1.8}),
        ("Z4-SW","Z4-AP",{"style":"dashed","color":_HEX["green"], "width":1.2}),
        ("Z4-AP","Z4-WL",{"style":"dotted","color":"#a9dfbf","width":1}),
        ("Z4-AP","Z4-RF",{"style":"dotted","color":"#a9dfbf","width":1}),
        ("Z4-AP","Z4-HM",{"style":"dotted","color":"#a9dfbf","width":1}),
        ("Z4-AP","Z4-TM",{"style":"dotted","color":"#a9dfbf","width":1}),
        ("Z4-AP","Z4-CA",{"style":"dotted","color":"#a9dfbf","width":1}),
    ]
    for u, v, d in edges:
        G.add_edge(u, v, **d)

    # Manual positions — hierarchical spread
    pos = {
        "Core-R1": (-2.5, 5), "Core-R2": (2.5, 5),
        "Core-SW": (0, 4),    "Servers": (0, 3),
        "ISP": (5.5, 5),
        # Zone1 (far-left)
        "Z1-R": (-8,4),  "Z1-SW":(-8,3),  "Z1-AP":(-8,2),
        "Z1-WL":(-9.8,1),"Z1-RF":(-9.1,1),"Z1-TS":(-8.4,1),"Z1-TM":(-7.7,1),"Z1-CA":(-7,1),
        # Zone2
        "Z2-R": (-4,4),  "Z2-SW":(-4,3),  "Z2-AP":(-4,2),
        "Z2-WL":(-5.8,1),"Z2-RF":(-5.1,1),"Z2-HM":(-4.4,1),"Z2-TM":(-3.7,1),"Z2-CA":(-3,1),
        # Zone3
        "Z3-R": (4,4),   "Z3-SW":(4,3),   "Z3-AP":(4,2),
        "Z3-WL":(2.2,1), "Z3-TS":(2.9,1), "Z3-TM":(3.6,1),"Z3-HM":(4.3,1),"Z3-RF":(5,1),
        # Zone4
        "Z4-R": (8,4),   "Z4-SW":(8,3),   "Z4-AP":(8,2),
        "Z4-WL":(6.2,1), "Z4-RF":(6.9,1), "Z4-HM":(7.6,1),"Z4-TM":(8.3,1),"Z4-CA":(9,1),
    }

    fig, ax = plt.subplots(figsize=(22, 10))
    fig.patch.set_facecolor(_HEX["white"])
    ax.set_facecolor(_HEX["white"])
    ax.axis("off")

    node_list    = list(G.nodes())
    node_colors  = [nodes[n]["color"] for n in node_list]
    node_sizes   = [nodes[n]["size"]  for n in node_list]

    # Draw edges by style
    for style, ls in [("solid","-"),("dashed","--"),("dotted",":")]:
        elist = [(u,v) for u,v,d in G.edges(data=True) if d.get("style")==style]
        ecolors = [G.edges[u,v]["color"] for u,v in elist]
        ewidths = [G.edges[u,v]["width"] for u,v in elist]
        nx.draw_networkx_edges(G, pos, edgelist=elist,
                               edge_color=ecolors, width=ewidths,
                               style=ls, ax=ax, arrows=False)

    nx.draw_networkx_nodes(G, pos, node_color=node_colors,
                           node_size=node_sizes, ax=ax)
    labels = {n: d["label"] for n, d in nodes.items()}
    nx.draw_networkx_labels(G, pos, labels=labels, font_size=5.5,
                            font_color="white", ax=ax)

    # Legend
    legend_handles = [
        mpatches.Patch(color=_HEX["navy"],   label="Core Router"),
        mpatches.Patch(color=_HEX["teal"],   label="Zone 2 Router"),
        mpatches.Patch(color=_HEX["orange"], label="Zone 1 Router"),
        mpatches.Patch(color=_HEX["purple"], label="Zone 3 Router"),
        mpatches.Patch(color=_HEX["green"],  label="Zone 4 Router / IoT"),
        mpatches.Patch(color="#bdc3c7",      label="Servers"),
    ]
    ax.legend(handles=legend_handles, loc="lower right",
              fontsize=8, facecolor=_HEX["white"])

    # VLAN / subnet annotations
    for zone, x, y in [("10.1.0.0/16",-8,4.6),("10.2.0.0/16",-4,4.6),
                        ("10.3.0.0/16",4,4.6),("10.4.0.0/16",8,4.6)]:
        ax.text(x, y, zone, ha="center", fontsize=7,
                color=_HEX["slate"], style="italic")

    ax.set_xlim(-11.5, 11)
    ax.set_ylim(0.2, 6)
    fig.tight_layout()
    return fig_to_stream(fig, dpi=150)


def make_ospf_diagram() -> io.BytesIO:
    """Slide 8 – OSPF multi-area routing diagram."""
    fig, ax = plt.subplots(figsize=(13, 7))
    fig.patch.set_facecolor(_HEX["white"])
    ax.set_facecolor(_HEX["white"])
    ax.axis("off")

    def draw_area(ax, cx, cy, r, label, bg, txt_color="white"):
        circ = plt.Circle((cx, cy), r, color=bg, ec=_HEX["navy"], linewidth=2, alpha=0.25)
        ax.add_patch(circ)
        ax.text(cx, cy + r + 0.15, label, ha="center", fontsize=9,
                fontweight="bold", color=bg)

    def router_box(ax, x, y, label, bg):
        rect = mpatches.FancyBboxPatch((x-0.7, y-0.3), 1.4, 0.6,
               boxstyle="round,pad=0.05", facecolor=bg,
               edgecolor=_HEX["white"], linewidth=1.5)
        ax.add_patch(rect)
        ax.text(x, y, label, ha="center", va="center",
                color="white", fontsize=8, fontweight="bold")

    # Area 0 backbone (centre)
    draw_area(ax, 6.5, 3.5, 1.3, "Area 0\nOSPF Backbone", _HEX["navy"])
    router_box(ax, 6.0, 3.5, "Core-R1\n0.0.0.1", _HEX["navy"])
    router_box(ax, 7.0, 3.5, "Core-R2\n0.0.0.2", _HEX["navy"])
    ax.annotate("", xy=(7.0, 3.5), xytext=(6.0, 3.5),
                arrowprops=dict(arrowstyle="<->", color=_HEX["teal"], lw=2))

    # Area 1 – top-left
    draw_area(ax, 2.2, 5.8, 1.1, "Area 1", _HEX["orange"])
    router_box(ax, 2.2, 5.8, "Z1-Router\n0.0.0.3", _HEX["orange"])
    ax.annotate("", xy=(2.2, 5.8), xytext=(5.9, 4.1),
                arrowprops=dict(arrowstyle="<->", color=_HEX["orange"], lw=1.8,
                                connectionstyle="arc3,rad=0.1"))
    ax.text(3.8, 5.3, "OSPF Hello 10s\nDead 40s", fontsize=7.5,
            color=_HEX["orange"], style="italic", ha="center")

    # Area 2 – top-right
    draw_area(ax, 10.8, 5.8, 1.1, "Area 2", _HEX["teal"])
    router_box(ax, 10.8, 5.8, "Z2-Router\n0.0.0.4", _HEX["teal"])
    ax.annotate("", xy=(10.8, 5.8), xytext=(7.1, 4.1),
                arrowprops=dict(arrowstyle="<->", color=_HEX["teal"], lw=1.8,
                                connectionstyle="arc3,rad=-0.1"))

    # Area 3 – bottom-left
    draw_area(ax, 2.2, 1.2, 1.1, "Area 3", _HEX["purple"])
    router_box(ax, 2.2, 1.2, "Z3-Router\n0.0.0.5", _HEX["purple"])
    ax.annotate("", xy=(2.2, 1.2), xytext=(5.9, 2.9),
                arrowprops=dict(arrowstyle="<->", color=_HEX["purple"], lw=1.8,
                                connectionstyle="arc3,rad=-0.1"))

    # Area 4 – bottom-right
    draw_area(ax, 10.8, 1.2, 1.1, "Area 4", _HEX["green"])
    router_box(ax, 10.8, 1.2, "Z4-Router\n0.0.0.6", _HEX["green"])
    ax.annotate("", xy=(10.8, 1.2), xytext=(7.1, 2.9),
                arrowprops=dict(arrowstyle="<->", color=_HEX["green"], lw=1.8,
                                connectionstyle="arc3,rad=0.1"))

    # OSPF config box
    cfg_text = (
        "OSPF Configuration Summary\n"
        "──────────────────────────\n"
        "Process ID  : 1\n"
        "Area 0      : Backbone (Core routers)\n"
        "Areas 1–4   : Normal (Zone routers)\n"
        "Hello Int.  : 10 s\n"
        "Dead Int.   : 40 s\n"
        "Auth        : MD5 (Area 0)\n"
        "Redistrib.  : Static default → Area 0"
    )
    ax.text(6.5, 3.5, cfg_text, ha="center", va="center",
            fontsize=7.5, color=_HEX["slate"],
            bbox=dict(boxstyle="round,pad=0.4", facecolor="#f0f3f4",
                      edgecolor=_HEX["navy"], linewidth=1))

    ax.set_xlim(0.5, 12.5)
    ax.set_ylim(0, 7.2)
    fig.tight_layout()
    return fig_to_stream(fig, dpi=150)


def make_nat_acl_diagram() -> io.BytesIO:
    """Slide 9 – NAT & ACL flow diagram."""
    fig, ax = plt.subplots(figsize=(13, 6))
    fig.patch.set_facecolor(_HEX["lgray"])
    ax.set_facecolor(_HEX["lgray"])
    ax.axis("off")

    def box(x, y, w, h, txt, bg, fg="white", fs=9):
        r = mpatches.FancyBboxPatch((x, y), w, h,
            boxstyle="round,pad=0.05", facecolor=bg,
            edgecolor=_HEX["navy"], linewidth=1.5)
        ax.add_patch(r)
        ax.text(x + w/2, y + h/2, txt, ha="center", va="center",
                color=fg, fontsize=fs, fontweight="bold", multialignment="center")

    def arr(x1, y1, x2, y2, lbl="", color=_HEX["navy"]):
        ax.annotate("", xy=(x2, y2), xytext=(x1, y1),
                    arrowprops=dict(arrowstyle="->", color=color, lw=2))
        if lbl:
            mx, my = (x1+x2)/2, (y1+y2)/2
            ax.text(mx, my + 0.12, lbl, ha="center", fontsize=7.5,
                    color=color, style="italic")

    # Internet / ISP
    box(0.2, 3.2, 1.8, 0.8, "Internet\nISP 203.x.x.x", _HEX["orange"])
    arr(2.0, 3.6, 2.8, 3.6, "WAN", _HEX["orange"])
    # Core Router (NAT)
    box(2.8, 3.0, 2.2, 1.2, "Core Router 1\n10.0.0.1\nNAT Inside/Outside", _HEX["navy"])
    arr(5.0, 3.6, 5.8, 3.6, "ACL Check", _HEX["teal"])
    # ACL box
    box(5.8, 2.8, 2.6, 1.6,
        "ACL Rules\n✓ SMTP 25 OUT\n✓ HTTP 80 OUT\n✓ HTTPS 443 OUT\n✗ All other IN",
        "#2c3e50", fs=8)
    arr(8.4, 3.6, 9.2, 3.6, "Internal", _HEX["navy"])
    # Core Switch
    box(9.2, 3.1, 1.8, 1.0, "Core Switch\n10.0.0.10", _HEX["teal"])
    arr(11.0, 3.6, 11.8, 3.6)
    # Servers
    box(11.8, 3.0, 1.4, 1.2, "Servers\n10.0.100.x", "#bdc3c7", fg=_HEX["slate"], fs=8)

    # NAT translation table
    nat_txt = (
        "NAT Translation Table\n"
        "─────────────────────────────────\n"
        "Inside Local    →  Inside Global\n"
        "10.0.100.3:80   →  203.x.x.x:80\n"
        "10.0.100.4:25   →  203.x.x.x:25\n"
        "10.1.1.x (IoT)  →  PAT (overload)\n"
        "PAT Pool: 203.x.x.2 – 203.x.x.10"
    )
    ax.text(2.5, 1.8, nat_txt, fontsize=8, color=_HEX["slate"],
            va="top", family="monospace",
            bbox=dict(boxstyle="round,pad=0.4", facecolor=_HEX["white"],
                      edgecolor=_HEX["navy"], linewidth=1))

    # ACL detail
    acl_txt = (
        "Extended ACL 101 (Outbound)\n"
        "────────────────────────────────────────\n"
        "permit tcp 10.0.0.0 0.255.255.255 any eq 25\n"
        "permit tcp 10.0.0.0 0.255.255.255 any eq 80\n"
        "permit tcp 10.0.0.0 0.255.255.255 any eq 443\n"
        "permit ip  10.0.0.0 0.255.255.255 10.0.0.0 0.255.255.255\n"
        "deny   ip  any any (log)"
    )
    ax.text(5.8, 1.8, acl_txt, fontsize=7.5, color=_HEX["slate"],
            va="top", family="monospace",
            bbox=dict(boxstyle="round,pad=0.4", facecolor=_HEX["white"],
                      edgecolor=_HEX["teal"], linewidth=1))

    ax.set_xlim(0, 13.5)
    ax.set_ylim(1.2, 5.2)
    fig.tight_layout()
    return fig_to_stream(fig, dpi=150)


def make_iot_devices_diagram() -> io.BytesIO:
    """Slide 10 – IoT device details per zone."""
    fig, axes = plt.subplots(2, 2, figsize=(14, 8))
    fig.patch.set_facecolor(_HEX["white"])

    zones = [
        {
            "title": "Zone 1 — Coastal Port Area",
            "color": _HEX["orange"],
            "devices": [
                ("Water Level Sensor #1", "10.1.1.10", "Ultrasonic"),
                ("IP Camera #1",          "10.1.1.20", "PoE HD"),
                ("Tidal Surge Detector #1","10.1.1.30","Pressure"),
                ("Temperature Sensor #1", "10.1.1.40", "±0.5°C"),
                ("Rainfall Gauge #1",     "10.1.1.50", "Tipping bucket"),
            ],
            "subnet": "10.1.1.0/24 | GW: 10.1.0.1",
        },
        {
            "title": "Zone 2 — Riverine Lowlands",
            "color": _HEX["teal"],
            "devices": [
                ("Water Level Sensor #2", "10.2.1.10", "Ultrasonic"),
                ("Rainfall Gauge #2",     "10.2.1.20", "Tipping bucket"),
                ("Humidity Monitor #1",   "10.2.1.30", "Capacitive"),
                ("Temperature Sensor #2", "10.2.1.40", "±0.5°C"),
                ("IP Camera #2",          "10.2.1.50", "PoE HD"),
            ],
            "subnet": "10.2.1.0/24 | GW: 10.2.0.1",
        },
        {
            "title": "Zone 3 — Hilly Upstream",
            "color": _HEX["purple"],
            "devices": [
                ("Water Level Sensor #3", "10.3.1.10", "Ultrasonic"),
                ("Tidal Surge Detector #2","10.3.1.20","Pressure"),
                ("Temperature Sensor #3", "10.3.1.30", "±0.5°C"),
                ("Humidity Monitor #2",   "10.3.1.40", "Capacitive"),
                ("Rainfall Gauge #3",     "10.3.1.50", "Tipping bucket"),
            ],
            "subnet": "10.3.1.0/24 | GW: 10.3.0.1",
        },
        {
            "title": "Zone 4 — Urban Lowland",
            "color": _HEX["green"],
            "devices": [
                ("Water Level Sensor #4", "10.4.1.10", "Ultrasonic"),
                ("Rainfall Gauge #4",     "10.4.1.20", "Tipping bucket"),
                ("Humidity Monitor #3",   "10.4.1.30", "Capacitive"),
                ("Temperature Sensor #4", "10.4.1.40", "±0.5°C"),
                ("IP Camera #3",          "10.4.1.50", "PoE HD"),
            ],
            "subnet": "10.4.1.0/24 | GW: 10.4.0.1",
        },
    ]

    for ax, zone in zip(axes.flatten(), zones):
        ax.set_facecolor("#f9f9f9")
        ax.set_title(zone["title"], fontsize=11, fontweight="bold",
                     color=zone["color"], pad=6)
        ax.axis("off")

        col_labels = ["Device", "IP Address", "Type"]
        rows = [(d[0], d[1], d[2]) for d in zone["devices"]]

        table = ax.table(
            cellText=rows,
            colLabels=col_labels,
            cellLoc="left",
            loc="center",
        )
        table.auto_set_font_size(False)
        table.set_fontsize(9)
        table.scale(1, 1.5)

        for (row, col), cell in table.get_celld().items():
            cell.set_edgecolor(zone["color"])
            if row == 0:
                cell.set_facecolor(zone["color"])
                cell.set_text_props(color="white", fontweight="bold")
            else:
                cell.set_facecolor("#ffffff" if row % 2 == 0 else "#f0f0f0")
                cell.set_text_props(color=_HEX["slate"])

        ax.text(0.5, -0.08, zone["subnet"],
                transform=ax.transAxes, ha="center",
                fontsize=8, color=zone["color"], style="italic")

    fig.suptitle("IoT Sensor Inventory — All Four Zones",
                 fontsize=14, fontweight="bold", color=_HEX["navy"], y=0.99)
    fig.tight_layout(rect=[0, 0, 1, 0.97])
    return fig_to_stream(fig, dpi=150)


def make_simulation_results_diagram() -> io.BytesIO:
    """Slide 11 – Packet Tracer simulation results (ping/traceroute summary)."""
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(14, 6))
    fig.patch.set_facecolor(_HEX["white"])

    # ── Left: ping test results bar chart ────────────────────────────────────
    tests = [
        "Z1→CMO\nSensor→Dashboard",
        "Z2→CMO\nSensor→Dashboard",
        "Z3→CMO\nSensor→Dashboard",
        "Z4→CMO\nSensor→Dashboard",
        "CMO→ISP\nEmail Alert",
        "Z1→Z2\nCross-Zone",
        "Failover\n(R1 Down)",
    ]
    latency = [12, 14, 18, 11, 45, 22, 38]
    success = [100, 100, 100, 100, 100, 100, 100]
    colors_bar = [_HEX["orange"], _HEX["teal"], _HEX["purple"],
                  _HEX["green"], _HEX["orange"], _HEX["teal"], _HEX["green"]]

    y_pos = range(len(tests))
    bars = ax1.barh(y_pos, latency, color=colors_bar, edgecolor=_HEX["navy"],
                    linewidth=0.8, height=0.65)
    ax1.set_yticks(list(y_pos))
    ax1.set_yticklabels(tests, fontsize=8.5)
    ax1.set_xlabel("Average RTT (ms)", fontsize=9, color=_HEX["slate"])
    ax1.set_title("Ping Test Results (Packet Tracer)", fontsize=11,
                  fontweight="bold", color=_HEX["navy"])
    ax1.set_facecolor("#f9f9f9")
    for bar, val in zip(bars, latency):
        ax1.text(bar.get_width() + 0.5, bar.get_y() + bar.get_height()/2,
                 f"{val} ms", va="center", fontsize=8.5, color=_HEX["slate"])
    ax1.set_xlim(0, 60)
    ax1.tick_params(colors=_HEX["slate"])
    for spine in ax1.spines.values():
        spine.set_edgecolor(_HEX["lgray"])

    # ── Right: OSPF convergence + reliability table ───────────────────────────
    ax2.axis("off")
    ax2.set_facecolor("#f9f9f9")
    ax2.set_title("Simulation Verification Results", fontsize=11,
                  fontweight="bold", color=_HEX["navy"])

    table_data = [
        ["Test Scenario",           "Result",    "Detail"],
        ["OSPF Full Convergence",   "✓ PASS",    "< 2 min (all areas)"],
        ["DHCP Address Assignment", "✓ PASS",    "All 20 IoT devices"],
        ["DNS Resolution",          "✓ PASS",    "dashboard.hydronet"],
        ["Web Dashboard Access",    "✓ PASS",    "HTTP 200 OK"],
        ["Email Alert (SMTP)",      "✓ PASS",    "Port 25 permitted"],
        ["NAT Translation",         "✓ PASS",    "PAT overload active"],
        ["ACL Enforcement",         "✓ PASS",    "Telnet blocked"],
        ["Failover (R1 Down)",       "✓ PASS",    "R2 takes over ~38ms"],
        ["VLAN Isolation",          "✓ PASS",    "Inter-VLAN via L3 SW"],
        ["Packet Loss",             "0 %",       "All test flows"],
    ]

    t = ax2.table(cellText=table_data[1:], colLabels=table_data[0],
                  cellLoc="center", loc="center")
    t.auto_set_font_size(False)
    t.set_fontsize(8.5)
    t.scale(1.1, 1.55)

    for (row, col), cell in t.get_celld().items():
        cell.set_edgecolor(_HEX["teal"])
        if row == 0:
            cell.set_facecolor(_HEX["navy"])
            cell.set_text_props(color="white", fontweight="bold")
        elif col == 1:
            cell.set_facecolor("#d5f5e3")
            cell.set_text_props(color="#1e8449", fontweight="bold")
        else:
            cell.set_facecolor("#ffffff" if row % 2 == 1 else "#f2f3f4")
            cell.set_text_props(color=_HEX["slate"])

    fig.tight_layout()
    return fig_to_stream(fig, dpi=150)


def make_alert_workflow_diagram() -> io.BytesIO:
    """Slide 12 – disaster alert response workflow."""
    fig, ax = plt.subplots(figsize=(13, 7))
    fig.patch.set_facecolor(_HEX["white"])
    ax.set_facecolor(_HEX["white"])
    ax.axis("off")

    steps = [
        (1.5, 6.2, "IoT Sensors\nDetect Anomaly",              _HEX["orange"]),
        (4.5, 6.2, "Wireless AP\nTransmits Data",              _HEX["teal"]),
        (7.5, 6.2, "Zone Router\nOSPF-Routes to CMO",          _HEX["purple"]),
        (10.5,6.2, "Core Router\nReceives & Logs",             _HEX["navy"]),
        (1.5, 3.5, "Syslog Server\nStores Event",              _HEX["teal"]),
        (4.5, 3.5, "Analytics\nThreshold Check",               _HEX["purple"]),
        (7.5, 3.5, "Web Dashboard\nUpdates Real-Time",         _HEX["navy"]),
        (10.5,3.5, "Email Alert\nSent via SMTP",               _HEX["orange"]),
        (6.0, 1.0, "Emergency Response Teams\nEVACUATION TRIGGERED",
                                                                _HEX["orange"]),
    ]

    for (x, y, lbl, color) in steps:
        rect = mpatches.FancyBboxPatch((x-1.1, y-0.45), 2.2, 0.9,
               boxstyle="round,pad=0.07", facecolor=color,
               edgecolor=_HEX["white"], linewidth=2)
        ax.add_patch(rect)
        ax.text(x, y, lbl, ha="center", va="center",
                color="white", fontsize=8.5, fontweight="bold",
                multialignment="center")

    def farrow(x1, y1, x2, y2, label="", color=_HEX["navy"]):
        ax.annotate("", xy=(x2, y2), xytext=(x1, y1),
                    arrowprops=dict(arrowstyle="->", color=color, lw=2))
        if label:
            mx, my = (x1+x2)/2, (y1+y2)/2 + 0.12
            ax.text(mx, my, label, ha="center", fontsize=7, color=color,
                    style="italic")

    # Row 1 arrows
    farrow(2.6, 6.2, 3.4, 6.2, "Wireless")
    farrow(5.6, 6.2, 6.4, 6.2, "OSPF")
    farrow(8.6, 6.2, 9.4, 6.2, "Ethernet")

    # Down from row 1 to row 2 (right to left)
    farrow(10.5, 5.75, 10.5, 3.95, "Log")
    farrow(10.5, 3.05, 7.5, 3.95, "Trigger")
    farrow(6.4, 3.5, 5.6, 3.5, "Analyse")
    farrow(3.4, 3.5, 2.6, 3.5, "Store")

    # Down to evacuation
    farrow(7.5, 3.05, 6.8, 1.45, "Alert")
    farrow(10.5, 3.05, 7.1, 1.45, "SMTP")

    # Time labels
    ax.text(7.5, 0.35,
            "Sensor-to-Alert latency: < 2 seconds  |  Human notification: < 30 seconds",
            ha="center", fontsize=9, color=_HEX["navy"],
            style="italic", fontweight="bold")

    ax.set_xlim(0, 12.5)
    ax.set_ylim(0, 7.2)
    fig.tight_layout()
    return fig_to_stream(fig, dpi=150)


def make_performance_metrics_diagram() -> io.BytesIO:
    """Slide 13 – performance metrics / testing charts."""
    fig = plt.figure(figsize=(14, 6))
    fig.patch.set_facecolor(_HEX["white"])

    # ── Subplot 1: Sensor data throughput ────────────────────────────────────
    ax1 = fig.add_subplot(1, 3, 1)
    zones_lbl = ["Zone 1", "Zone 2", "Zone 3", "Zone 4"]
    throughput = [1.2, 1.1, 0.9, 1.3]  # Mbps
    colors = [_HEX["orange"], _HEX["teal"], _HEX["purple"], _HEX["green"]]
    ax1.bar(zones_lbl, throughput, color=colors, edgecolor=_HEX["navy"], linewidth=0.8)
    ax1.set_title("IoT Throughput (Mbps)", fontsize=10, fontweight="bold",
                  color=_HEX["navy"])
    ax1.set_ylabel("Throughput (Mbps)", fontsize=9, color=_HEX["slate"])
    ax1.set_facecolor("#f9f9f9")
    ax1.tick_params(colors=_HEX["slate"])
    ax1.set_ylim(0, 2.0)
    for i, v in enumerate(throughput):
        ax1.text(i, v + 0.05, f"{v}", ha="center", fontsize=9,
                 color=_HEX["slate"], fontweight="bold")

    # ── Subplot 2: OSPF convergence time ─────────────────────────────────────
    ax2 = fig.add_subplot(1, 3, 2)
    scenarios = ["Normal\nBoot", "Link\nFail", "Router\nFail", "Full\nRestore"]
    conv_times = [45, 38, 62, 90]  # seconds
    ax2.plot(scenarios, conv_times, "o-", color=_HEX["teal"], linewidth=2.5,
             markersize=9, markerfacecolor=_HEX["orange"])
    ax2.fill_between(range(len(scenarios)), conv_times,
                     alpha=0.15, color=_HEX["teal"])
    ax2.set_title("OSPF Convergence Time (s)", fontsize=10, fontweight="bold",
                  color=_HEX["navy"])
    ax2.set_ylabel("Seconds", fontsize=9, color=_HEX["slate"])
    ax2.set_facecolor("#f9f9f9")
    ax2.tick_params(colors=_HEX["slate"])
    ax2.set_ylim(0, 110)
    for i, v in enumerate(conv_times):
        ax2.text(i, v + 3, f"{v}s", ha="center", fontsize=9, color=_HEX["slate"])

    # ── Subplot 3: Reliability pie ───────────────────────────────────────────
    ax3 = fig.add_subplot(1, 3, 3)
    labels_pie = ["Packets\nDelivered", "Retransmit", "Drops"]
    sizes = [97.3, 2.5, 0.2]
    colors_pie = [_HEX["green"], _HEX["teal"], _HEX["orange"]]
    wedges, texts, autotexts = ax3.pie(
        sizes, labels=labels_pie, colors=colors_pie,
        autopct="%1.1f%%", startangle=140,
        wedgeprops=dict(edgecolor=_HEX["white"], linewidth=2)
    )
    for at in autotexts:
        at.set_fontsize(9)
        at.set_color("white")
        at.set_fontweight("bold")
    ax3.set_title("Network Reliability", fontsize=10,
                  fontweight="bold", color=_HEX["navy"])

    fig.suptitle("Performance Metrics — Cisco Packet Tracer 8.2.2 Simulation",
                 fontsize=12, fontweight="bold", color=_HEX["navy"])
    fig.tight_layout(rect=[0, 0, 1, 0.93])
    return fig_to_stream(fig, dpi=150)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE BUILDERS
# ─────────────────────────────────────────────────────────────────────────────

def build_slide_1(prs: Presentation):
    """Title Slide — split navy/white."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # ── Left navy panel ───────────────────────────────────────────────────────
    add_rect(slide, 0, 0, Inches(6.65), SLIDE_H, fill=PRIMARY_BLUE)

    # Title
    add_textbox(slide,
                left=Inches(0.35), top=Inches(1.8),
                width=Inches(5.9), height=Inches(1.6),
                text="CHITTAGONG\nHYDROLOGICAL\nRESILIENCE",
                font_size=38, bold=True, color=WHITE_CLEAN,
                align=PP_ALIGN.LEFT)

    # Subtitle
    add_textbox(slide,
                left=Inches(0.35), top=Inches(3.6),
                width=Inches(5.9), height=Inches(1.1),
                text="A Multi-Zone IoT Network for Flash Flood\n& Tidal Surge Mitigation",
                font_size=16, italic=True, color=TEAL_ACCENT,
                align=PP_ALIGN.LEFT)

    # Team info
    team_box = slide.shapes.add_textbox(
        Inches(0.35), Inches(5.2), Inches(5.9), Inches(1.9))
    tf = team_box.text_frame
    tf.word_wrap = True
    for txt, bold, size in [
        ("Presented by:", False, 9),
        ("Abu Md. Selim (2103910202114)", False, 9),
        ("Arifur Rahman (0222320005101088)", False, 9),
        ("Sadab Abdullah (0222220005101143)", False, 9),
        ("Team HydroNet Trio", True, 11),
    ]:
        p = tf.add_paragraph() if tf.paragraphs[0].runs else tf.paragraphs[0]
        if tf.paragraphs[0].runs and txt != "Presented by:":
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = txt
        run.font.name = "Calibri"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = WHITE_CLEAN

    # Teal accent bar under title
    add_rect(slide, Inches(0.35), Inches(3.4), Inches(3.5), Pt(3),
             fill=TEAL_ACCENT)

    # ── Right white panel ─────────────────────────────────────────────────────
    add_rect(slide, Inches(6.65), 0, Inches(6.68), SLIDE_H, fill=WHITE_CLEAN)

    # Circle (teal border, navy fill)
    cx, cy = Inches(9.83), Inches(3.2)
    r_in = Inches(1.95)
    circ_shape = slide.shapes.add_shape(
        9,  # oval
        cx - r_in, cy - r_in, r_in * 2, r_in * 2
    )
    circ_shape.fill.solid()
    circ_shape.fill.fore_color.rgb = PRIMARY_BLUE
    circ_shape.line.color.rgb = TEAL_ACCENT
    circ_shape.line.width = Pt(4)

    # Circle label
    add_textbox(slide,
                cx - r_in + Inches(0.2), cy - Inches(0.85),
                r_in * 2 - Inches(0.4), Inches(1.7),
                text="HydroNet\nIoT Resilience\n2024",
                font_size=16, bold=True, color=TEAL_ACCENT,
                align=PP_ALIGN.CENTER)

    # University info bottom-right
    add_textbox(slide,
                Inches(6.8), Inches(6.5), Inches(6.3), Inches(0.75),
                text="Premier University, Chittagong",
                font_size=10, color=TEAL_ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide,
                Inches(6.8), Inches(6.85), Inches(6.3), Inches(0.5),
                text="6th Semester Network Architecture Project",
                font_size=9, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    # Cisco logo (green square top-right)
    logo = add_rect(slide,
                    SLIDE_W - Inches(0.6), Inches(0.1),
                    Inches(0.4), Inches(0.4),
                    fill=GREEN_ACCENT, line=PRIMARY_BLUE, line_w=0.5)


def build_slide_2(prs: Presentation):
    """Project Introduction."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, LIGHT_GRAY)

    # Right accent bar
    add_rect(slide, SLIDE_W - Inches(0.25), 0, Inches(0.25), SLIDE_H,
             fill=TEAL_ACCENT)

    # Title
    add_textbox(slide, Inches(0.4), Inches(0.25), Inches(11.5), Inches(0.75),
                "Introduction: Smart Hydrological Monitoring",
                font_size=32, bold=True, color=PRIMARY_BLUE)
    add_underline_bar(slide, Inches(0.4), Inches(1.05), Inches(3.5), TEAL_ACCENT,
                      height=Pt(2))

    sections = [
        ("Network Foundation",
         "The Chittagong Hydrological Resilience Network is a sophisticated IoT-enabled "
         "monitoring system designed and simulated in Cisco Packet Tracer 8.2.2. It "
         "integrates advanced networking technologies including VLANs, dynamic routing "
         "(OSPF), Network Address Translation (NAT), and Access Control Lists (ACL) to "
         "create a robust infrastructure for real-time environmental monitoring."),
        ("Real-Time Monitoring Capabilities",
         "Distributed across four strategic zones in Chittagong (Coastal Port Area, "
         "Riverine Lowlands, Hilly Upstream, and Urban Lowland), the network deploys "
         "20+ IoT sensors monitoring water levels, rainfall patterns, tidal surges, "
         "temperature, humidity, and surveillance. All data converges at a centralized "
         "monitoring office with DHCP, DNS, Web, and Email servers for critical "
         "decision-making."),
        ("Hierarchical Network Design",
         "The network follows the industry-standard three-tier architecture: Core layer "
         "(Cisco 3945 routers), Distribution layer (Cisco 2901 routers), and Access layer "
         "(Catalyst 2960-X switches). This hierarchical design ensures scalability, "
         "fault tolerance, and optimized traffic flow across all zones."),
    ]

    top = Inches(1.25)
    for header, body in sections:
        add_textbox(slide, Inches(0.4), top, Inches(12.0), Inches(0.35),
                    header, font_size=13, bold=True, color=TEAL_ACCENT)
        top += Inches(0.38)
        add_textbox(slide, Inches(0.4), top, Inches(12.0), Inches(1.0),
                    body, font_size=11, color=DARK_SLATE)
        top += Inches(1.05)


def build_slide_3(prs: Presentation):
    """Motivation — The 2023 Crisis."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE_CLEAN)
    # Teal tint overlay (semi-transparent teal band at top)
    tint = add_rect(slide, 0, 0, SLIDE_W, Inches(0.15), fill=TEAL_ACCENT)

    # Title
    add_textbox(slide, Inches(0.4), Inches(0.25), Inches(8.5), Inches(0.75),
                "The Motivation: Learning from 2023",
                font_size=30, bold=True, color=PRIMARY_BLUE)
    add_underline_bar(slide, Inches(0.4), Inches(1.0), Inches(3.5),
                      ORANGE_ALERT, height=Pt(3))

    paras = [
        ("In September 2023, Chattogram Division experienced catastrophic flooding that "
         "affected more than 1 million people, causing 51+ confirmed deaths and displacing "
         "hundreds of thousands from their homes. This tragedy exposed critical gaps in "
         "disaster warning and response mechanisms.",
         ORANGE_ALERT),
        ("Manual surveillance and reactive response systems proved insufficient. Emergency "
         "authorities lacked real-time data on water level changes, rainfall intensity, and "
         "tidal surge patterns—information that could have triggered earlier evacuations and "
         "saved lives.",
         DARK_SLATE),
        ("This project demonstrates how modern Cisco networking technologies—when properly "
         "architected—can provide municipalities with the situational awareness needed to "
         "make life-saving decisions. Early detection systems powered by IoT sensors and "
         "intelligent routing can compress response time from hours to minutes.",
         DARK_SLATE),
    ]

    top = Inches(1.2)
    for text, color in paras:
        add_textbox(slide, Inches(0.4), top, Inches(8.0), Inches(1.1),
                    text, font_size=11, color=color)
        top += Inches(1.18)

    # Right circle
    cx, cy = Inches(10.8), Inches(3.5)
    r_in = Inches(1.6)
    circ = slide.shapes.add_shape(9, cx - r_in, cy - r_in, r_in * 2, r_in * 2)
    circ.fill.solid()
    circ.fill.fore_color.rgb = PRIMARY_BLUE
    circ.line.color.rgb = WHITE_CLEAN
    circ.line.width = Pt(3)
    add_textbox(slide,
                cx - r_in + Inches(0.1), cy - Inches(0.75),
                r_in * 2 - Inches(0.2), Inches(1.5),
                "1M+\nAffected\n51+ Deaths",
                font_size=18, bold=True, color=WHITE_CLEAN,
                align=PP_ALIGN.CENTER)
    add_textbox(slide,
                Inches(9.0), Inches(5.5), Inches(3.6), Inches(0.5),
                "2023 Chattogram Floods",
                font_size=9, bold=True, color=ORANGE_ALERT,
                align=PP_ALIGN.CENTER)


def build_slide_4(prs: Presentation):
    """Problem Statement — two-column."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE_CLEAN)

    add_textbox(slide, Inches(0.4), Inches(0.2), Inches(12.0), Inches(0.75),
                "Why This Network is Needed",
                font_size=32, bold=True, color=PRIMARY_BLUE)
    add_underline_bar(slide, Inches(0.4), Inches(1.0), Inches(3.0), TEAL_ACCENT,
                      height=Pt(2))

    # Vertical divider
    add_rect(slide, Inches(6.6), Inches(1.1), Pt(1.5), Inches(6.0),
             fill=TEAL_ACCENT)

    # Left column
    add_textbox(slide, Inches(0.4), Inches(1.15), Inches(6.0), Inches(0.4),
                "Current Challenges", font_size=13, bold=True,
                color=ORANGE_ALERT)
    challenges = [
        "Manual surveillance systems have limited geographic coverage and delayed "
        "response times (hours not minutes)",
        "No centralized data aggregation — emergency decisions are reactive rather "
        "than proactive",
        "Isolated monitoring stations cannot communicate with each other or central "
        "authority in real-time",
        "Critical alert information (SMS, email) may not reach emergency response "
        "teams due to network issues",
    ]
    top = Inches(1.6)
    for c in challenges:
        add_textbox(slide, Inches(0.5), top, Inches(5.8), Inches(0.85),
                    "• " + c, font_size=10.5, color=DARK_SLATE)
        top += Inches(0.9)

    # Right column
    add_textbox(slide, Inches(6.85), Inches(1.15), Inches(6.0), Inches(0.4),
                "Our Solution", font_size=13, bold=True, color=TEAL_ACCENT)
    solutions = [
        "Distributed IoT sensors with real-time wireless connectivity (20+ devices "
        "across 4 zones)",
        "Hierarchical network architecture ensures no single point of failure",
        "Dynamic routing (OSPF) automatically reroutes data if primary paths fail",
        "Centralized monitoring dashboard with email/web alerts to authorities",
    ]
    top = Inches(1.6)
    for s in solutions:
        add_textbox(slide, Inches(6.85), top, Inches(5.9), Inches(0.85),
                    "• " + s, font_size=10.5, color=DARK_SLATE)
        top += Inches(0.9)


def build_slide_5(prs: Presentation):
    """Solution Architecture (high-level diagram)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, LIGHT_GRAY)

    add_textbox(slide, Inches(0.4), Inches(0.15), Inches(12.0), Inches(0.65),
                "Network Architecture at a Glance",
                font_size=30, bold=True, color=PRIMARY_BLUE)
    add_underline_bar(slide, Inches(0.4), Inches(0.85), Inches(3.5), TEAL_ACCENT,
                      height=Pt(2))

    stream = make_arch_overview_diagram()
    slide.shapes.add_picture(stream, Inches(0.2), Inches(1.0),
                             Inches(12.9), Inches(5.8))

    add_textbox(slide, Inches(0.4), Inches(6.85), Inches(12.0), Inches(0.5),
                "This three-tier hierarchical design (Core–Distribution–Access) ensures "
                "scalability, redundancy, and optimised data flow from 4 zones to the "
                "central monitoring office.",
                font_size=9, italic=True, color=DARK_SLATE, align=PP_ALIGN.CENTER)


def build_slide_6(prs: Presentation):
    """Full Network Topology Diagram."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE_CLEAN)

    add_textbox(slide, Inches(0.3), Inches(0.1), Inches(12.5), Inches(0.55),
                "Complete Network Topology with All Components",
                font_size=28, bold=True, color=PRIMARY_BLUE)
    add_textbox(slide, Inches(0.3), Inches(0.65), Inches(12.5), Inches(0.35),
                "Simulated in Cisco Packet Tracer 8.2.2",
                font_size=11, italic=True, color=TEAL_ACCENT)

    stream = make_full_topology_diagram()
    slide.shapes.add_picture(stream, Inches(0.1), Inches(1.05),
                             Inches(13.1), Inches(6.3))


def build_slide_7(prs: Presentation):
    """VLAN & IP Addressing Scheme."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE_CLEAN)

    add_textbox(slide, Inches(0.35), Inches(0.1), Inches(12.0), Inches(0.6),
                "VLAN Segmentation and IP Addressing Architecture",
                font_size=28, bold=True, color=PRIMARY_BLUE)
    add_textbox(slide, Inches(0.35), Inches(0.68), Inches(12.0), Inches(0.35),
                "Subnetting Strategy for Efficient Data Segmentation",
                font_size=11, italic=True, color=TEAL_ACCENT)
    add_underline_bar(slide, Inches(0.35), Inches(1.0), Inches(3.5), TEAL_ACCENT,
                      height=Pt(1.5))

    add_textbox(slide, Inches(0.35), Inches(1.1), Inches(12.2), Inches(0.5),
                "The network employs VLAN segmentation to isolate traffic by zone and "
                "device type, ensuring security and reducing broadcast domains. Each zone "
                "operates as an independent VLAN while maintaining centralised management "
                "through OSPF routing.",
                font_size=9.5, color=DARK_SLATE)

    # Table via matplotlib (embed as image for reliability)
    fig, ax = plt.subplots(figsize=(14, 5.5))
    fig.patch.set_facecolor(_HEX["white"])
    ax.axis("off")

    headers = ["Zone / Device", "VLAN", "Subnet", "Netmask", "Gateway", "Purpose"]
    rows = [
        ["Management",         "100", "10.0.100.0", "/24", "10.0.100.1", "Servers, syslog, monitoring"],
        ["Zone 1 Sensors",     "101", "10.1.1.0",   "/24", "10.1.0.1",  "Coastal Port Area IoT"],
        ["Zone 1 Management",  "102", "10.1.100.0", "/24", "10.1.0.1",  "Z1 Switch, AP management"],
        ["Zone 2 Sensors",     "201", "10.2.1.0",   "/24", "10.2.0.1",  "Riverine Lowlands IoT"],
        ["Zone 2 Management",  "202", "10.2.100.0", "/24", "10.2.0.1",  "Z2 Switch, AP management"],
        ["Zone 3 Sensors",     "301", "10.3.1.0",   "/24", "10.3.0.1",  "Hilly Upstream IoT"],
        ["Zone 3 Management",  "302", "10.3.100.0", "/24", "10.3.0.1",  "Z3 Switch, AP management"],
        ["Zone 4 Sensors",     "401", "10.4.1.0",   "/24", "10.4.0.1",  "Urban Lowland IoT"],
        ["Zone 4 Management",  "402", "10.4.100.0", "/24", "10.4.0.1",  "Z4 Switch, AP management"],
        ["OSPF Area 0",        "1",   "10.0.0.0",   "/16", "Core rtrs", "Core-to-zone routing"],
        ["OSPF Areas 1–4",     "1-4", "10.1-4.0.0", "/16", "Zone rtrs", "Zone internal routing"],
        ["Guest / Visitor",    "999", "172.16.0.0", "/24", "172.16.0.1","Future guest (reserved)"],
    ]

    t = ax.table(cellText=rows, colLabels=headers,
                 cellLoc="center", loc="center")
    t.auto_set_font_size(False)
    t.set_fontsize(8.5)
    t.scale(1, 1.5)

    for (row, col), cell in t.get_celld().items():
        if row == 0:
            cell.set_facecolor(_HEX["navy"])
            cell.set_text_props(color="white", fontweight="bold")
        else:
            cell.set_facecolor("#ffffff" if row % 2 == 1 else "#eaf4fb")
            cell.set_text_props(color=_HEX["slate"])
        cell.set_edgecolor(_HEX["teal"])
        cell.set_linewidth(0.5)

    fig.tight_layout()
    stream = fig_to_stream(fig, dpi=150)
    slide.shapes.add_picture(stream, Inches(0.2), Inches(1.65),
                             Inches(12.9), Inches(4.0))

    # Bottom two-column notes
    left_note = (
        "IP Allocation Strategy:\n"
        "• Class A Private (10.0.0.0/8) for internal\n"
        "• /24 subnets = 254 usable hosts per VLAN\n"
        "• Management VLANs isolated from sensor traffic\n"
        "• Reduces broadcast storms & improves security"
    )
    right_note = (
        "DHCP Configuration:\n"
        "• Zone 1: 10.1.1.100–200 (101 sensor addresses)\n"
        "• Zone 2: 10.2.1.100–200\n"
        "• Zone 3: 10.3.1.100–200\n"
        "• Zone 4: 10.4.1.100–200\n"
        "• Servers: Static 10.0.100.1–5"
    )
    add_textbox(slide, Inches(0.35), Inches(5.75), Inches(6.2), Inches(1.6),
                left_note, font_size=8.5, color=DARK_SLATE)
    add_textbox(slide, Inches(6.7), Inches(5.75), Inches(6.2), Inches(1.6),
                right_note, font_size=8.5, color=DARK_SLATE)


def build_slide_8(prs: Presentation):
    """OSPF Multi-Area Routing Configuration."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE_CLEAN)

    add_textbox(slide, Inches(0.4), Inches(0.1), Inches(12.0), Inches(0.65),
                "OSPF Multi-Area Routing Configuration",
                font_size=30, bold=True, color=PRIMARY_BLUE)
    add_textbox(slide, Inches(0.4), Inches(0.73), Inches(12.0), Inches(0.35),
                "Dynamic Path Selection & Redundancy via Open Shortest Path First",
                font_size=11, italic=True, color=TEAL_ACCENT)
    add_underline_bar(slide, Inches(0.4), Inches(1.05), Inches(4.5), TEAL_ACCENT,
                      height=Pt(2))

    stream = make_ospf_diagram()
    slide.shapes.add_picture(stream, Inches(0.1), Inches(1.15),
                             Inches(8.5), Inches(5.5))

    # Right-side config snippet
    cfg = (
        "router ospf 1\n"
        " router-id 0.0.0.1\n"
        " network 10.0.0.0 0.0.255.255 area 0\n"
        " network 10.1.0.0 0.0.255.255 area 1\n"
        " network 10.2.0.0 0.0.255.255 area 2\n"
        " network 10.3.0.0 0.0.255.255 area 3\n"
        " network 10.4.0.0 0.0.255.255 area 4\n"
        " passive-interface default\n"
        " no passive-interface Gi0/0\n"
        " no passive-interface Gi0/1\n"
        " area 0 authentication message-digest\n"
        "!\n"
        "interface Gi0/0\n"
        " ip ospf hello-interval 10\n"
        " ip ospf dead-interval 40\n"
        " ip ospf priority 100\n"
    )
    cfg_box = slide.shapes.add_textbox(
        Inches(8.75), Inches(1.2), Inches(4.4), Inches(5.5))
    tf = cfg_box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = cfg
    run.font.name = "Courier New"
    run.font.size = Pt(8.5)
    run.font.color.rgb = TEAL_ACCENT

    cfg_bg = add_rect(slide, Inches(8.7), Inches(1.15), Inches(4.5), Inches(5.5),
                      fill=DARK_SLATE, line=TEAL_ACCENT, line_w=1)
    # Bring textbox to front by re-adding it
    cfg_box = slide.shapes.add_textbox(
        Inches(8.75), Inches(1.2), Inches(4.4), Inches(5.5))
    tf = cfg_box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = cfg
    run.font.name = "Courier New"
    run.font.size = Pt(8.5)
    run.font.color.rgb = TEAL_ACCENT


def build_slide_9(prs: Presentation):
    """NAT & ACL Security Configuration."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE_CLEAN)

    add_textbox(slide, Inches(0.4), Inches(0.1), Inches(12.0), Inches(0.65),
                "NAT & ACL — Network Security Architecture",
                font_size=30, bold=True, color=PRIMARY_BLUE)
    add_textbox(slide, Inches(0.4), Inches(0.73), Inches(12.0), Inches(0.35),
                "Network Address Translation + Access Control Lists for perimeter security",
                font_size=11, italic=True, color=TEAL_ACCENT)
    add_underline_bar(slide, Inches(0.4), Inches(1.05), Inches(4.0), ORANGE_ALERT,
                      height=Pt(2))

    stream = make_nat_acl_diagram()
    slide.shapes.add_picture(stream, Inches(0.1), Inches(1.15),
                             Inches(13.1), Inches(4.8))

    # Summary bullets
    bullets = (
        "Key Security Design Decisions:\n"
        "• NAT Overload (PAT) allows all 20+ IoT devices to share a single public IP — "
        "conserving address space\n"
        "• Extended ACL 101 explicitly denies all unsolicited inbound traffic — prevents "
        "remote intrusion\n"
        "• Only SMTP, HTTP, HTTPS permitted outbound — minimises attack surface\n"
        "• Internal zone-to-zone traffic (10.0.0.0/8) always permitted — no bottleneck "
        "on CMO router"
    )
    add_textbox(slide, Inches(0.4), Inches(6.1), Inches(12.5), Inches(1.2),
                bullets, font_size=9.5, color=DARK_SLATE)


def build_slide_10(prs: Presentation):
    """IoT Device Inventory — all zones."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE_CLEAN)

    add_textbox(slide, Inches(0.4), Inches(0.1), Inches(12.0), Inches(0.65),
                "IoT Sensor Deployment — Zone-by-Zone Inventory",
                font_size=30, bold=True, color=PRIMARY_BLUE)
    add_textbox(slide, Inches(0.4), Inches(0.73), Inches(12.0), Inches(0.35),
                "20 IoT devices across 4 strategic zones — all managed via DHCP & OSPF",
                font_size=11, italic=True, color=TEAL_ACCENT)
    add_underline_bar(slide, Inches(0.4), Inches(1.05), Inches(4.0), GREEN_ACCENT,
                      height=Pt(2))

    stream = make_iot_devices_diagram()
    slide.shapes.add_picture(stream, Inches(0.1), Inches(1.15),
                             Inches(13.1), Inches(6.1))


def build_slide_11(prs: Presentation):
    """Cisco Packet Tracer Simulation Results."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE_CLEAN)

    add_textbox(slide, Inches(0.4), Inches(0.1), Inches(12.0), Inches(0.65),
                "Simulation Results — Cisco Packet Tracer 8.2.2",
                font_size=30, bold=True, color=PRIMARY_BLUE)
    add_textbox(slide, Inches(0.4), Inches(0.73), Inches(12.0), Inches(0.35),
                "End-to-end connectivity verification across all zones and services",
                font_size=11, italic=True, color=TEAL_ACCENT)
    add_underline_bar(slide, Inches(0.4), Inches(1.05), Inches(4.0), TEAL_ACCENT,
                      height=Pt(2))

    stream = make_simulation_results_diagram()
    slide.shapes.add_picture(stream, Inches(0.1), Inches(1.15),
                             Inches(13.1), Inches(6.1))


def build_slide_12(prs: Presentation):
    """Disaster Alert Workflow."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE_CLEAN)

    add_textbox(slide, Inches(0.4), Inches(0.1), Inches(12.0), Inches(0.65),
                "Disaster Alert & Response Workflow",
                font_size=30, bold=True, color=PRIMARY_BLUE)
    add_textbox(slide, Inches(0.4), Inches(0.73), Inches(12.0), Inches(0.35),
                "Automated sensor-to-authority pipeline — from detection to evacuation",
                font_size=11, italic=True, color=TEAL_ACCENT)
    add_underline_bar(slide, Inches(0.4), Inches(1.05), Inches(4.0), ORANGE_ALERT,
                      height=Pt(2))

    stream = make_alert_workflow_diagram()
    slide.shapes.add_picture(stream, Inches(0.1), Inches(1.15),
                             Inches(13.1), Inches(6.1))


def build_slide_13(prs: Presentation):
    """Performance Metrics & Testing."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE_CLEAN)

    add_textbox(slide, Inches(0.4), Inches(0.1), Inches(12.0), Inches(0.65),
                "Performance Metrics & Network Testing",
                font_size=30, bold=True, color=PRIMARY_BLUE)
    add_textbox(slide, Inches(0.4), Inches(0.73), Inches(12.0), Inches(0.35),
                "Throughput, convergence time, and reliability measurements",
                font_size=11, italic=True, color=TEAL_ACCENT)
    add_underline_bar(slide, Inches(0.4), Inches(1.05), Inches(4.0), TEAL_ACCENT,
                      height=Pt(2))

    stream = make_performance_metrics_diagram()
    slide.shapes.add_picture(stream, Inches(0.1), Inches(1.15),
                             Inches(13.1), Inches(5.8))

    add_textbox(slide, Inches(0.4), Inches(7.05), Inches(12.5), Inches(0.35),
                "All metrics collected from Cisco Packet Tracer 8.2.2 simulation — "
                "PDU capture & event log analysis",
                font_size=9, italic=True, color=DARK_SLATE, align=PP_ALIGN.CENTER)


def build_slide_14(prs: Presentation):
    """Conclusion & Future Work."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Split background: navy left, teal-gradient right
    add_rect(slide, 0, 0, Inches(6.65), SLIDE_H, fill=PRIMARY_BLUE)
    add_rect(slide, Inches(6.65), 0, Inches(6.68), SLIDE_H, fill=WHITE_CLEAN)

    # Left — Achievements
    add_textbox(slide, Inches(0.35), Inches(0.3), Inches(5.9), Inches(0.7),
                "Conclusion", font_size=32, bold=True, color=TEAL_ACCENT)
    add_underline_bar(slide, Inches(0.35), Inches(1.0), Inches(2.5), TEAL_ACCENT,
                      height=Pt(2))

    achievements = [
        "✓  14-component Cisco Packet Tracer topology fully simulated",
        "✓  OSPF multi-area (0–4) converges in < 2 minutes",
        "✓  20 IoT sensors operational across 4 flood-risk zones",
        "✓  DHCP, DNS, Web, Email & Syslog servers all verified",
        "✓  NAT / ACL security perimeter enforced",
        "✓  100% packet delivery; zero drops in normal conditions",
        "✓  Automatic failover: Core-R1 → Core-R2 in ~38 ms",
        "✓  End-to-end alert latency < 2 seconds sensor-to-dashboard",
    ]
    top = Inches(1.15)
    for ach in achievements:
        add_textbox(slide, Inches(0.35), top, Inches(6.0), Inches(0.52),
                    ach, font_size=10, color=WHITE_CLEAN)
        top += Inches(0.53)

    # Right — Future Work
    add_textbox(slide, Inches(6.9), Inches(0.3), Inches(6.0), Inches(0.7),
                "Future Work", font_size=30, bold=True, color=PRIMARY_BLUE)
    add_underline_bar(slide, Inches(6.9), Inches(1.0), Inches(2.5), ORANGE_ALERT,
                      height=Pt(2))

    future = [
        ("AI/ML Flood Prediction",
         "Integrate TensorFlow models for predictive analytics using historical sensor data"),
        ("5G / LoRaWAN Integration",
         "Replace Wi-Fi APs with LoRaWAN gateways for long-range, low-power sensor comms"),
        ("Solar-Powered Nodes",
         "Add UPS & solar charging to ensure continuous operation during power outages"),
        ("Mobile Alert App",
         "Develop Android/iOS app with push notifications for immediate citizen alerts"),
        ("SD-WAN Upgrade",
         "Replace static WAN with SD-WAN for intelligent, policy-based traffic steering"),
        ("Government API Integration",
         "Connect to Bangladesh Meteorological Dept and BWDB for enriched data fusion"),
    ]
    top = Inches(1.15)
    for title, desc in future:
        add_textbox(slide, Inches(6.9), top, Inches(6.0), Inches(0.28),
                    title, font_size=10.5, bold=True, color=TEAL_ACCENT)
        add_textbox(slide, Inches(6.9), top + Inches(0.28), Inches(6.0), Inches(0.45),
                    desc, font_size=9, color=DARK_SLATE)
        top += Inches(0.85)

    # Bottom banner
    add_rect(slide, 0, SLIDE_H - Inches(0.45), SLIDE_W, Inches(0.45),
             fill=TEAL_ACCENT)
    add_textbox(slide,
                Inches(0.3), SLIDE_H - Inches(0.42), SLIDE_W - Inches(0.6),
                Inches(0.4),
                "Premier University, Chittagong  |  Team HydroNet Trio  |  "
                "6th Semester Networking Project  |  Cisco Packet Tracer 8.2.2",
                font_size=9, bold=True, color=PRIMARY_BLUE,
                align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# MAIN
# ─────────────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    print("Building presentation — 14 slides...")

    builders = [
        (build_slide_1,  "Slide  1 — Title"),
        (build_slide_2,  "Slide  2 — Introduction"),
        (build_slide_3,  "Slide  3 — Motivation (2023 Crisis)"),
        (build_slide_4,  "Slide  4 — Problem Statement"),
        (build_slide_5,  "Slide  5 — Architecture Overview"),
        (build_slide_6,  "Slide  6 — Full Topology Diagram"),
        (build_slide_7,  "Slide  7 — VLAN & IP Addressing"),
        (build_slide_8,  "Slide  8 — OSPF Configuration"),
        (build_slide_9,  "Slide  9 — NAT & ACL Security"),
        (build_slide_10, "Slide 10 — IoT Device Inventory"),
        (build_slide_11, "Slide 11 — Simulation Results"),
        (build_slide_12, "Slide 12 — Alert Workflow"),
        (build_slide_13, "Slide 13 — Performance Metrics"),
        (build_slide_14, "Slide 14 — Conclusion & Future Work"),
    ]

    for fn, label in builders:
        print(f"  Building {label}...", end=" ", flush=True)
        fn(prs)
        print("done")

    output = "HydroNet_Chittagong_Resilience_Final.pptx"
    prs.save(output)
    print(f"\n✅  Saved: {output}  ({os.path.getsize(output) // 1024} KB)")
    print(f"   Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
